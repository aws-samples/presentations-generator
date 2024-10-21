## Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
## SPDX-License-Identifier: MIT-0

import boto3
import os
import io
import pptx
from pptx.util import Inches
import logging
from boto3.dynamodb.conditions import Key

# Configurar clientes AWS
dynamodb = boto3.resource('dynamodb')
s3 = boto3.client('s3')

# Configurar parâmetros
table_name = os.environ["table_name"]
bucket_name = os.environ["bucket_name"]
path = os.environ["path"]
template_key = os.environ["template_key"]

def lambda_handler(event, context):
    logging.info('Função Lambda iniciada')

    # Obter dados do DynamoDB
    session_id = event['session_id']
    logging.info(f'Obtendo dados do DynamoDB para session_id: {session_id}')
    dynamodb_data = get_dynamodb_data(table_name, session_id)
    logging.info(f'Dados obtidos do DynamoDB: {dynamodb_data}')

    
    # Obter URLs das imagens no S3
    prefix_imagens = f'{path}/{session_id}/'
    imagem_urls = get_imagem_urls(bucket_name, prefix_imagens)

    # Obter o arquivo de template do bucket S3
    template_file = io.BytesIO()
    s3.download_fileobj(bucket_name, template_key, template_file)
    template_file.seek(0)

    # Criar apresentação PowerPoint
    presentation = pptx.Presentation(template_file)

    # Adicionar slides com textos e imagens
    for data, imagem_key in zip(dynamodb_data, imagem_urls):
        print("Dados sendo processados:", data)
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Slide com título e conteúdo
        if slide.shapes.title:
            slide.shapes.title.text = data['sldtitle']
        else:
            print(f"Erro: slide.shapes.title é None para o item: {data}")
        body_shape = slide.shapes.placeholders[1]
        body_shape.text = f"{data['sldcontent']}"

        # Encontrar o espaço reservado para imagens
        picture_holder_index = []
        for shape in slide.shapes:
                if shape.is_placeholder:
                    if "PICTURE" in str(shape.placeholder_format.type):
                        file_obj = download_file(bucket_name, imagem_key)
                        picture_holder_index.append(shape.placeholder_format.idx)
                        # insert graph
                        slide.placeholders[picture_holder_index[0]].insert_picture(file_obj)
                else:
                    print(f"Erro: não foi possível encontrar o espaço reservado para imagens no slide para o item: {data}")

    # Salvar apresentação no S3
    output_prefix = f'{path}/{session_id}/'  # Defina o prefixo desejado
    output_filename = f'{output_prefix}apresentacao.pptx'
    file_object = io.BytesIO()
    presentation.save(file_object)
    file_object.seek(0)

    s3.upload_fileobj(file_object, bucket_name, output_filename)

    # Atualizar status do DynamoDB
    table = dynamodb.Table(table_name)
    response = table.update_item(
        Key={'session_id': session_id},
        UpdateExpression="SET output_filename = :output_filename",
        ExpressionAttributeValues={':output_filename': output_filename},
        ReturnValues="UPDATED_NEW"
    )
    print(f"Item atualizado: {response['Attributes']}")

    return {
        'statusCode': 200,
        'body': f'A apresentação "{output_filename}" foi salva no bucket "{bucket_name}".'
    }

def get_dynamodb_data(table_name, session_id):
    table = dynamodb.Table(table_name)

    # Filtra os itens pelo valor do campo 'session_id'
    response = table.query(
        KeyConditionExpression=Key('session_id').eq(session_id)
    )

    data = response['Items']

    # Extrai os valores dos campos 'sldtitle' e 'sldcontent' dos slides
    slides = []
    for item in data:
        for slide in item['slides']:
            slide_data = {
                'sldtitle': slide['sldtitle'],
                'sldcontent': slide['sldcontent']
            }
            slides.append(slide_data)

    return slides

def get_imagem_urls(bucket_name, prefix):
    s3_client = boto3.client('s3')
    response = s3_client.list_objects_v2(Bucket=bucket_name, Prefix=prefix)
    imagem_keys = []
    for obj in response.get('Contents', []):
        objeto_key = obj['Key']
        if objeto_key != prefix:  # Ignorar a chave que representa o prefixo
            imagem_keys.append(objeto_key)
            print(f"Chave do objeto: {objeto_key}")  # Adicione este print
        else:
            print(f"A chave do objeto '{objeto_key}' não começa com o prefixo '{prefix}'")

    return imagem_keys

def download_file(bucket_name, object_key):
    file_obj = io.BytesIO()
    s3.download_fileobj(bucket_name, object_key, file_obj)
    file_obj.seek(0)
    return file_obj

if __name__ == "__main__":
    context = None
    try:
        result = lambda_handler({}, context)
        print(result)
    except Exception as e:
        logging.error(f"Ocorreu uma exceção: {e}")